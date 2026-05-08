from pptx import Presentation
prs = Presentation("Akshayam Tejomaya.pptx")
slide = prs.slides[9] # Slide 10
for shape in slide.shapes:
    if shape.has_table:
        header_row = shape.table.rows[0]
        header_texts = [cell.text.strip() for cell in header_row.cells]
        print(header_texts)
