from pptx import Presentation

prs = Presentation("Akshayam Tejomaya.pptx")
for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        if shape.has_table:
            header_row = shape.table.rows[0]
            header_texts = [cell.text.strip() for cell in header_row.cells]
            normalized = [" ".join(h.split()) for h in header_texts]
            print(f"Slide {i+1}, Shape {j}: {normalized}")
