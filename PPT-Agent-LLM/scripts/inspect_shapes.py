from pptx import Presentation

def inspect_slide_shapes(ppt_path, slide_index):
    prs = Presentation(ppt_path)
    if slide_index < len(prs.slides):
        slide = prs.slides[slide_index]
        print(f"--- Shapes on Slide {slide_index + 1} ---")
        for i, shape in enumerate(slide.shapes):
            print(f"Shape {i}: Name='{shape.name}', ID={shape.shape_id}, Type={shape.shape_type}")
            if hasattr(shape, "text"):
                print(f"  Text: {shape.text[:100]}...")

if __name__ == "__main__":
    ppt_path = "Akshayam Tejomaya edition No 407 Week 46 10.11.2025 to 16.11.2025 (1) (1).pptx"
    inspect_slide_shapes(ppt_path, 4) # Slide 5 is index 4
