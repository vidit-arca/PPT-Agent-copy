import glob
import os
from pptx import Presentation

# Find the most recently modified PPT file that looks like a generated one
ppt_files = glob.glob("tejomaya_filled_*.pptx")
if not ppt_files:
    print("No generated PPT files found.")
    exit()

latest_ppt = max(ppt_files, key=os.path.getctime)
print(f"Verifying: {latest_ppt}")

prs = Presentation(latest_ppt)

# Check if slide index 3 exists
if len(prs.slides) > 3:
    slide = prs.slides[3]  # Slide 4 (0-indexed)
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            print(f"Table has {len(table.rows)} rows")
            
            for row_idx in range(min(6, len(table.rows))):  # First 6 rows
                row = table.rows[row_idx]
                print(f"\nRow {row_idx}:")
                for col_idx, cell in enumerate(row.cells):
                    print(f"  Col {col_idx}: '{cell.text[:50]}'")
else:
    print("Slide 4 not found in the presentation.")
