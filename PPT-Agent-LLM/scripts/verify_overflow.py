import pandas as pd
from pptx import Presentation
from agent import fill_using_mapping

def verify():
    ppt_path = "test_overflow.pptx"
    excel_path = "test_data.xlsx"
    output_path = "verified_overflow.pptx"

    print(f"Loading {ppt_path}...")
    prs = Presentation(ppt_path)
    
    print(f"Loading {excel_path}...")
    df = pd.read_excel(excel_path)
    
    # Construct manual mapping for the test
    # We target Slide 0, Shape 1 (Table 1)
    mapping = {
        "tables": [
            {
                "slide_index": 0,
                "shape_index": 1,
                "headers": {
                    "S. No": None,
                    "Date of Issue": ["IssueDate"],
                    "Rules / circulars / Notifications / Order": ["SubCategory"],
                    "Contents thereof": ["Title"],
                    "Gist thereof": ["Summary"]
                },
                "source_filter": {"Verticals": "Test Vertical"},
                "max_rows": 100
            }
        ]
    }
    
    print("Running fill_using_mapping...")
    updates = fill_using_mapping(prs, df, mapping)
    
    print(f"Updates made: {updates}")
    
    prs.save(output_path)
    print(f"Saved to {output_path}")
    
    # Verification checks
    prs_out = Presentation(output_path)
    print(f"Output slides: {len(prs_out.slides)}")
    
    if len(prs_out.slides) >= 2:
        print("✓ Slide count increased (Expected >= 2)")
        
        slide1 = prs_out.slides[0]
        slide2 = prs_out.slides[1]
        
        # Check Slide 1
        # Should have Table 1 (filled) and NO Table 2
        # Table 1 is shape 1. Table 2 was shape 2.
        # If Table 2 was removed, shape count should be less?
        # Or indices shifted.
        print(f"Slide 1 shapes: {len(slide1.shapes)}")
        for i, s in enumerate(slide1.shapes):
            if s.has_table:
                tbl = s.table
                print(f"  Shape {i} (Table): {len(tbl.rows)} rows. Cell(0,0): {tbl.cell(0,0).text}")
        
        # Check Slide 2
        # Should have Table 2 (empty/template) and NO Table 1
        print(f"Slide 2 shapes: {len(slide2.shapes)}")
        for i, s in enumerate(slide2.shapes):
            if s.has_table:
                tbl = s.table
                print(f"  Shape {i} (Table): {len(tbl.rows)} rows. Cell(0,0): {tbl.cell(0,0).text}")

    else:
        print("✗ Slide count did NOT increase!")

if __name__ == "__main__":
    verify()
