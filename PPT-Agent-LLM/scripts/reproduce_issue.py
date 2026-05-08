import os
from pptx import Presentation
from pptx.util import Inches
from agent import fill_using_mapping
import pandas as pd

def create_dummy_image():
    from PIL import Image
    img = Image.new('RGB', (100, 100), color = 'red')
    img.save('test_image.png')

def reproduce():
    create_dummy_image()
    
    prs = Presentation()
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add an image (this creates a relationship)
    slide.shapes.add_picture('test_image.png', Inches(0.5), Inches(0.5), width=Inches(1))
    
    # Add Table 1
    rows, cols = 2, 5
    left, top, width, height = Inches(1), Inches(2), Inches(8), Inches(1)
    table_shape1 = slide.shapes.add_table(rows, cols, left, top, width, height)
    table1 = table_shape1.table
    # Set headers
    headers = ["S. No", "Date of Issue", "Rules / circulars / Notifications / Order", "Contents thereof", "Gist thereof"]
    for i, h in enumerate(headers):
        table1.cell(0, i).text = h

    # Add Table 2 (below)
    top2 = Inches(5)
    table_shape2 = slide.shapes.add_table(rows, cols, left, top2, width, height)
    table2 = table_shape2.table
    for i, h in enumerate(headers):
        table2.cell(0, i).text = f"Table 2 - {h}"
        
    prs.save("test_issue.pptx")
    
    # Create data to trigger overflow
    data = {
        "Verticals": ["Test Vertical"] * 20,
        "IssueDate": [f"2023-01-{i:02d}" for i in range(1, 21)],
        "SubCategory": [f"Rule {i}" for i in range(1, 21)],
        "Title": [f"Title {i}" for i in range(1, 21)],
        "Summary": [f"Summary {i}" for i in range(1, 21)]
    }
    df = pd.DataFrame(data)
    
    # Mapping
    mapping = {
        "tables": [
            {
                "slide_index": 0,
                "shape_index": 1, # Table 1 (Shape 0 is image)
                "headers": {
                    "S. No": None,
                    "Date of Issue": ["IssueDate"],
                    "Rules / circulars / Notifications / Order": ["SubCategory"],
                    "Contents thereof": ["Title"],
                    "Gist thereof": ["Summary"]
                },
                "source_filter": {"Verticals": "Test Vertical"},
                "max_rows": 100
            }
        ]
    }
    
    print("Running fill_using_mapping...")
    fill_using_mapping(prs, df, mapping)
    
    output_path = "verified_issue.pptx"
    prs.save(output_path)
    print(f"Saved to {output_path}")
    
    # Verify
    prs_out = Presentation(output_path)
    if len(prs_out.slides) >= 2:
        print("✓ Slide count increased")
        slide2 = prs_out.slides[1]
        
        # Check for image
        has_image = False
        for s in slide2.shapes:
            if s.shape_type == 13: # PICTURE
                has_image = True
                print("✓ Image found on Slide 2")
                break
        if not has_image:
            print("✗ Image MISSING on Slide 2")
            
        # Check Table 2 position
        # Table 2 should be near the top (where Table 1 was, ~2 inches)
        # Table 1 was at 2.0 inches.
        # Table 2 was at 5.0 inches.
        # Shift should be ~3.0 inches.
        # So Table 2 should be at ~2.0 inches.
        
        table2_found = False
        for s in slide2.shapes:
            if s.has_table and "Table 2" in s.table.cell(0,0).text:
                print(f"✓ Table 2 found at top: {s.top/914400:.2f} inches")
                if s.top < Inches(3):
                    print("✓ Table 2 is positioned correctly (near top)")
                else:
                    print(f"✗ Table 2 is still low: {s.top/914400:.2f} inches")
                table2_found = True
                break
        
        if not table2_found:
            print("✗ Table 2 NOT found on Slide 2")

if __name__ == "__main__":
    reproduce()
