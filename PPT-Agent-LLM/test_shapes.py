from pptx import Presentation

prs = Presentation("tejomaya_filled_20260822_122640.pptx")
slide = prs.slides[4] # Slide 5 is SEBI (index 4)
for i, shape in enumerate(slide.shapes):
    print(f"Shape {i}: Type {shape.shape_type}, Name: {shape.name}, Top: {shape.top}, Height: {shape.height}")
