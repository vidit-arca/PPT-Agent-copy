
import logging
from pptx import Presentation
from pptx.util import Inches

def inspect_slide9():
    ppt_path = "Akshayam Tejomaya.pptx"
    prs = Presentation(ppt_path)
    
    slide = prs.slides[8]
    print(f"Slide Height: {prs.slide_height.inches:.2f} inches")
    
    print(f"--- Inspecting Slide 9 ({slide.name}) ---")
    
    # Find the main table
    table_shape = None
    for shape in slide.shapes:
        if shape.has_table:
            table_shape = shape
            break
            
    if not table_shape:
        print("No table found!")
        return

    print(f"Main Table Top: {table_shape.top.inches:.2f} inches")
    print(f"Main Table Height: {table_shape.height.inches:.2f} inches")
    table_bottom_emu = table_shape.top + table_shape.height
    print(f"Main Table Bottom: {table_bottom_emu / 914400:.2f} inches")
    
    # Find shapes below table header (approx 0.5 inch)
    threshold = table_shape.top + Inches(0.5)
    
    print(f"\nListing ALL shapes on Slide 9:")
    
    for shape in slide.shapes:
        print(f"  Shape: ID={shape.shape_id}, Name='{shape.name}', Top={shape.top.inches:.2f}, Height={shape.height.inches:.2f}")
        if shape.has_text_frame:
            print(f"    Text: '{shape.text_frame.text}'")
        if shape.has_table:
            print(f"    Table with {len(shape.table.rows)} rows")
            
    if footer_shapes:
        min_top = min(s.top for s in footer_shapes)
        max_bot = max(s.top + s.height for s in footer_shapes)
        height_emu = max_bot - min_top
        height_inches = height_emu.inches
        print(f"\nCalculated Footer Height: {height_inches:.2f} inches")
        print(f"Footer Top: {min_top.inches:.2f} inches")
        print(f"Footer Bottom: {max_bot.inches:.2f} inches")
    else:
        print("\nNo footer shapes found.")

if __name__ == "__main__":
    inspect_slide9()
