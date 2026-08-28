import json
import os
import argparse
import glob
import logging
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches

from src.services.ppt.template import extract_template_spec
from src.services.data.profiler import build_data_profile
from src.config.settings import get_config
from src.services.llm.provider import get_llm_provider, ask_llm_for_mapping_with_retry
from src.utils.logging import setup_logging
from src.services.ppt.engine import (
    duplicate_slide, 
    add_table_rows, 
    remove_table_rows,
    get_shapes_below, 
    set_cell_text, 
    update_slide_subtitle,
    remove_shape,
    format_date
)
from src.services.ppt.layout import LayoutEngine

from src.services.text.processor import AITextProcessor
from src.services.storage.minio_client import MinioClient

# MANUAL VERTICAL-TO-SLIDE MAPPING# Manual slide mapping: Vertical → Slide Index (0-indexed)
MANUAL_SLIDE_MAPPING = {
    "Companies Act": 3,                   # Slide 4 - "Companies Act 2013 - MCA"
    "SEBI": 4,                           # Slide 5 - "SEBI during the week"
    "RBI": 5,                            # Slide 6 - "RBI during the week"
    "IFSC": 6,                           # Slide 7 - "IFSC during the week"
    "Insolvency and Bankruptcy Code": 7, # Slide 8 - "Insolvency and Bankruptcy Code"
    "Alternate Investment Funds": 8,      # Slide 9 - "Alternate Investment Funds"
    "Others during the week": 9,          # Slide 10 - "Others during the week"
}

# Subdomain-specific slide mapping: (Subdomain, Vertical) → Slide Index
# This allows different subdomains to go to different slides
# NOTE: Leave empty to combine all subdomains of the same vertical into one table
# If you want to separate subdomains into different slides, add entries here
SUBDOMAIN_SLIDE_MAPPING = {
    # Example: ("Consultation Paper", "SEBI"): 5,  # Would put Consultation Paper on Slide 6
    # Currently empty = all SEBI subdomains go to the same table
}

# FILE-TO-VERTICAL MAPPING
# Maps Excel file names (without extension) to their vertical categories
# This allows multiple Excel files to be processed, each representing a different vertical
# If a file is not in this mapping, the file name itself will be used as the vertical
FILE_VERTICAL_MAPPING = {
    "AIF": "Alternate Investment Funds",
    "Listed Companies": "Others during the week",
    "SEBI": "SEBI",
    "IBBI": "Insolvency and Bankruptcy Code",
    "Companies Act": "Companies Act",
    "RBI": "RBI",
    "ICAI": "Others during the week",
    # Add more mappings as needed
}


def get_llm():
    """Get LLM provider (Phase 1: Enhanced with OpenAI support and fallback)."""
    config = get_config()
    return get_llm_provider()


def ask_llm_for_mapping(llm, template_spec: dict, data_profile: dict) -> dict:
    """Ask LLM to generate DYNAMIC mapping (Phase 1: Enhanced with structured outputs and retry).
    
    NOTE: Temporarily disabled - AI mapping needs prompt improvements.
    Using manual fallback for reliability.
    """
    logging.info("AI mapping temporarily disabled, using manual fallback for reliability")
    return get_fallback_mapping(template_spec, data_profile)


def get_fallback_mapping(template_spec: dict, data_profile: dict) -> dict:
    """Fallback: use manual slide mapping and first 5-col table on each slide.
    Now supports subdomain-specific mappings.
    """
    logging.info("=" * 70)
    logging.info("FALLBACK MAPPING (Using Manual Slide Indexes)")
    logging.info("=" * 70)

    unique_verticals = data_profile.get("unique_verticals", [])
    unique_subdomains = data_profile.get("unique_subdomains", [])
    
    logging.info(f"Verticals in Excel: {unique_verticals}")
    if unique_subdomains:
        logging.info(f"Subdomains in Excel: {unique_subdomains}")
    logging.info("Manual slide mapping:")

    for v, slide_idx in MANUAL_SLIDE_MAPPING.items():
        logging.debug(f"  {v} → Slide {slide_idx + 1}")

    mappings = []

    # Create mappings for each subdomain-vertical combination
    if unique_subdomains:
        for subdomain in unique_subdomains:
            for vertical in unique_verticals:
                # Only create mapping if this subdomain actually exists for this vertical in the data
                # We can check this by filtering the data_profile['data'] if available, or just create it and let fill_using_mapping handle empty data
                # But to avoid "No mapping found" warnings for invalid combos, we should be smarter.
                # However, unique_subdomains is global.
                # Let's just create the mapping. fill_using_mapping will skip if no data matches.
                
                logging.info(f"Processing: Subdomain='{subdomain}', Vertical='{vertical}'")
                
                # Check subdomain-specific mapping first
                slide_idx = SUBDOMAIN_SLIDE_MAPPING.get((subdomain, vertical))
                
                # Fall back to vertical-only mapping
                if slide_idx is None:
                    slide_idx = MANUAL_SLIDE_MAPPING.get(vertical)
                
                # Partial match if exact key not present
                if slide_idx is None:
                    for key in MANUAL_SLIDE_MAPPING:
                        if key.lower() in vertical.lower() or vertical.lower() in key.lower():
                            slide_idx = MANUAL_SLIDE_MAPPING[key]
                            logging.info(f"  ✓ Partial match: '{vertical}' → '{key}' → Slide {slide_idx + 1}")
                            break

                if slide_idx is None:
                    # Don't warn here, as many combos might be invalid. 
                    # Only warn if we know for sure this combo exists in data.
                    # For now, just continue.
                    continue

                if slide_idx >= len(template_spec["slides"]):
                    logging.error(f"  ✗ Slide index {slide_idx} out of range")
                    continue

                slide_spec = template_spec["slides"][slide_idx]
                matched_shape = None
                for shape in slide_spec["shapes"]:
                    if shape["is_table"] and len(shape.get("table_headers", [])) == 5:
                        matched_shape = shape
                        break

                if not matched_shape:
                    logging.warning(f"  ✗ No 5-column table found on Slide {slide_idx + 1}")
                    continue

                # Create specific mapping for this subdomain
                source_filter = {"Verticals": vertical, "Subdomain": subdomain}
                logging.info(f"  → Mapping Subdomain='{subdomain}' for Vertical='{vertical}' to Slide {slide_idx + 1}")
                
                mappings.append({
                    "slide_index": slide_idx,
                    "shape_index": matched_shape["shape_index"],
                    "headers": {
                        "S. No": None,
                        "Date of Issue": ["IssueDate"],
                        "Rules / circulars / Notifications / Order": ["SubCategory", "PDF_URL"],
                        "Rules/ circulars / Notifications / Order": ["SubCategory", "PDF_URL"],
                        "Contents thereof": ["Title"],
                        "Gist thereof": ["Summary"],
                    },
                    "source_filter": source_filter,
                    "max_rows": 4,  # Limit to 4 rows per slide to prevent overflow
                })
    else:
        # Original logic for single-sheet files (no subdomains)
        for vertical in unique_verticals:
            logging.info(f"Processing vertical: '{vertical}'")
            slide_idx = MANUAL_SLIDE_MAPPING.get(vertical)

            # Partial match if exact key not present
            if slide_idx is None:
                for key in MANUAL_SLIDE_MAPPING:
                    if key.lower() in vertical.lower() or vertical.lower() in key.lower():
                        slide_idx = MANUAL_SLIDE_MAPPING[key]
                        logging.info(f"  ✓ Partial match: '{vertical}' → '{key}' → Slide {slide_idx + 1}")
                        break

            if slide_idx is None:
                logging.warning(f"  ✗ No mapping found for '{vertical}' – add it to MANUAL_SLIDE_MAPPING")
                continue

            if slide_idx >= len(template_spec["slides"]):
                logging.error(f"  ✗ Slide index {slide_idx} out of range for '{vertical}'")
                continue

            slide_spec = template_spec["slides"][slide_idx]
            matched_shape = None
            for shape in slide_spec["shapes"]:
                if shape["is_table"] and len(shape.get("table_headers", [])) == 5:
                    matched_shape = shape
                    logging.debug(f"  ✓ Found 5-column table at shape {shape['shape_index']}")
                    break

            if not matched_shape:
                logging.warning(f"  ✗ No 5-column table found on Slide {slide_idx + 1} for '{vertical}'")
                continue

            mappings.append({
                "slide_index": slide_idx,
                "shape_index": matched_shape["shape_index"],
                "headers": {
                    "S. No": None,
                    "Date of Issue": ["IssueDate"],
                    "Rules / circulars / Notifications / Order": ["SubCategory", "PDF_URL"],
                    "Rules/ circulars / Notifications / Order": ["SubCategory", "PDF_URL"],
                    "Contents thereof": ["Title"],
                    "Gist thereof": ["Summary"],
                },
                "source_filter": {"Verticals": vertical},
                "max_rows": 100,
            })
            logging.info(f"  ✓ MAPPED: '{vertical}' → Slide {slide_idx + 1}, Shape {matched_shape['shape_index']}")

    logging.info("=" * 70)
    logging.info(f"Total mappings created: {len(mappings)}")
    logging.info("=" * 70)

    return {"tables": mappings}


def paginate_and_fill(prs, slide_idx, shape_idx, df, headers_map, header_texts, header_to_colspec, precomputed_bullets_map, layout_engine=None, full_df=None, bottom_margin=0.5):
    """
    Recursively fill data into slides, creating new slides as needed.
    """
    if layout_engine is None:
        layout_engine = LayoutEngine()
        
    logging.info(f"--- Paginate and Fill: Slide {slide_idx+1}, Shape {shape_idx}, {len(df)} rows remaining ---")
    
    if len(df) == 0:
        return 0

    if slide_idx >= len(prs.slides):
        logging.error(f"✗ Slide index {slide_idx} out of range")
        return 0

    slide = prs.slides[slide_idx]
    
    # Update Subtitle if full_df is provided
    if full_df is not None:
        try:
            update_slide_subtitle(slide, full_df)
        except Exception as e:
            logging.warning(f"  ⚠ Failed to update subtitle: {e}")
    
    # Find the table shape
    shape = None
    if shape_idx < len(slide.shapes):
        shape = slide.shapes[shape_idx]
        # If the shape at this index is NOT a table, don't give up yet.
        # It might be that indices shifted or we were passed a wrong index.
        if not shape.has_table:
            logging.warning(f"  Shape at index {shape_idx} is not a table. Searching for table...")
            shape = None # Reset to trigger search below

    if not shape:
        # Try to find by proximity or just find the first table
        for s in slide.shapes:
            if s.has_table:
                shape = s
                break
    
    if not shape or not shape.has_table:
        logging.error(f"✗ Table shape not found on Slide {slide_idx+1}")
        return 0

    table = shape.table
    
    # --- Horizontal Overflow Check ---
    col_widths = [col.width.inches for col in table.columns]
    slide_width = prs.slide_width.inches
    layout_engine.check_horizontal_overflow(col_widths, slide_width)
    # ---------------------------------

    # --- Footer Height Calculation ---
    # Identify footer shapes (News table, etc.) to reserve space for them on the last slide
    # We assume anything below the table's initial bottom is part of the footer
    initial_table_bottom = shape.top + shape.height
    footer_shapes = get_shapes_below(slide, initial_table_bottom + Inches(0.5))
    # Filter out the table itself just in case
    footer_shapes = [s for s in footer_shapes if s.shape_id != shape.shape_id]
    
    footer_height_inches = 0.0
    if footer_shapes:
        try:
            min_top = min(s.top for s in footer_shapes)
            max_bot = max(s.top + s.height for s in footer_shapes)
            # Height of the footer block + spacing from table
            height_emu = max_bot - min_top
            footer_height_inches = (height_emu.inches if hasattr(height_emu, 'inches') else height_emu / 914400) + 0.3 # 0.3 inch spacing
            logging.info(f"  Footer height detected: {footer_height_inches:.2f} inches")
        except Exception as e:
            logging.warning(f"  Error calculating footer height: {e}")

    # Calculate available height
    slide_height_inches = prs.slide_height.inches
    
    # Initialize current_height from HEADER ONLY
    # We assume the table starts with just the header or we will clear it
    # But for estimation, we should start with the header height.
    # If the table has empty rows, we will remove them later.
    if len(table.rows) > 0:
        # Template XML header row can have an unrendered placeholder height (e.g. 1.90").
        # Real rendered header row with 1-2 lines of text is at most 0.5 - 0.6 inches.
        header_height = min(table.rows[0].height.inches, 0.6)
    else:
        header_height = 0.5 # Fallback
        
    current_height = header_height
    
    # Determine how many rows fit on THIS slide
    rows_for_this_slide = 0
    
    rows_fitting_texts = []
    
    # If this is a fresh overflow slide, we might have cleared rows but kept the header.
    # So current_table_height is just the header height.
    
    # Iterate through data to find split point
    for i, (_, row) in enumerate(df.iterrows()):
        # Check if this is a Section Header
        is_section_header = row.get('IsSectionHeader', False)
        
        # Build row text for estimation
        row_texts = []
        main_content = ""
        
        if is_section_header:
            # Section header usually spans, so it's one long text
            text_val = str(row.get('SectionTitle', 'Section'))
            row_texts = [text_val] # Treat as single column for height est (approx)
        else:
            for header_text in header_texts:
                colspec = header_to_colspec.get(header_text)
                if colspec is None: # S.No
                    row_texts.append("00")
                elif isinstance(colspec, list):
                    parts = []
                    for col_name in colspec:
                        if col_name in df.columns:
                            v = row[col_name]
                            if not pd.isna(v):
                                if col_name == "IssueDate":
                                    parts.append(format_date(v))
                                else:
                                    parts.append(str(v))
                    text_val = " ".join(parts)
                    
                    if "gist" in header_text.lower():
                        # Use precomputed bullets if available for more accurate layout estimation
                        bullets = precomputed_bullets_map.get(text_val)
                        if bullets:
                            text_val = "\n".join(bullets)
                            
                    row_texts.append(text_val)
                    
                    if "contents" in header_text.lower() or "gist" in header_text.lower():
                        if len(text_val) > len(main_content):
                            main_content = text_val
                else:
                    row_texts.append("")
        
        # Estimate height
        if is_section_header:
            # Section headers are usually single lines, maybe bold
            row_height = 0.4 # slightly larger than default
        else:
            row_height = layout_engine.estimate_row_height(row_texts, col_widths)
        
        # Check overflow
        # If this is the LAST row of the dataframe, we must also fit the footer!
        required_height = row_height
        # Footer reservation removed to allow table to fill slide.
        # Footer overflow will be handled at the end by moving it to a new slide.

        # Check overflow using dynamic height engine with 0.8" bottom margin safety
        table_top_inches = shape.top.inches
        effective_slide_height = slide_height_inches - table_top_inches
        
        if layout_engine.check_overflow(current_height, required_height, effective_slide_height, bottom_margin=bottom_margin):
            logging.info(f"  Overflow at row {i} (SectionHeader={is_section_header})")
            
            # If it's a section header, we definitely want to push it to the next slide
            # unless it's the VERY FIRST row, in which case we can't (infinite loop)
            if is_section_header and rows_for_this_slide > 0:
                break
                
            # Content-Aware Splitting for normal rows
            if not is_section_header and rows_fitting_texts:
                optimal_count = layout_engine.analyze_split_point(rows_fitting_texts, main_content)
                if optimal_count < rows_for_this_slide:
                    logging.info(f"  AI suggested backtracking split point from {rows_for_this_slide} to {optimal_count}")
                    rows_for_this_slide = optimal_count
            break
            
        current_height += row_height
        rows_for_this_slide += 1
        rows_fitting_texts.append(main_content)

    # Ensure at least 1 row if it's the very first row and it doesn't fit (to avoid infinite loop)
    if rows_for_this_slide == 0 and len(df) > 0:
        logging.warning("  Row 1 doesn't fit! Forcing it to avoid infinite loop.")
        rows_for_this_slide = 1

    logging.info(f"  Rows for this slide: {rows_for_this_slide}")

    # Slice data
    df_current = df.iloc[:rows_for_this_slide]
    df_remaining = df.iloc[rows_for_this_slide:]
    
    # --- FILL DATA ---
    # Add rows if needed (we assume table has header + maybe some rows)
    # We need to match table si
    # ze to df_current size
    # Current table rows (excluding header)
    current_data_rows = len(table.rows) - 1
    needed_rows = len(df_current)
    
    if needed_rows > current_data_rows:
        add_table_rows(table, needed_rows - current_data_rows)
    elif needed_rows < current_data_rows:
        # Remove extra rows to shrink table to fit content
        rows_to_remove = current_data_rows - needed_rows
        logging.info(f"  Removing {rows_to_remove} unused rows from table")
        remove_table_rows(table, rows_to_remove)

    # Fill rows
    row_idx = 1
    for i, (_, row) in enumerate(df_current.iterrows(), 1):
        is_section_header = row.get('IsSectionHeader', False)
        # Handle NaN (which is truthy) - ensure it's explicitly True
        if pd.isna(is_section_header):
            is_section_header = False
        
        if is_section_header:
            # Merge cells for section header
            # We want to merge all columns in this row
            first_cell = table.cell(row_idx, 0)
            other_cell = table.cell(row_idx, len(table.columns) - 1)
            first_cell.merge(other_cell)
            
            # Set text
            title = str(row.get('SectionTitle', ''))
            set_cell_text(first_cell, title)
            logging.info(f"    Filled Section Header: {title}")
            
        else:
            # Normal Row
            for col_idx, header_text in enumerate(header_texts):
                colspec = header_to_colspec.get(header_text)
                
                # S. No
                if colspec is None:
                    if header_text.lower().startswith("s. no"):
                        set_cell_text(table.cell(row_idx, col_idx), str(row.name + 1 if isinstance(row.name, int) else row_idx))
                    continue
                
                cell = table.cell(row_idx, col_idx)
                if isinstance(colspec, list):
                    parts = []
                    for col_name in colspec:
                        if col_name in df.columns:
                            v = row[col_name]
                            if not pd.isna(v):
                                if col_name == "IssueDate":
                                    parts.append(format_date(v))
                                else:
                                    parts.append(str(v))
                    
                    # Use newline for Rules column to ensure paragraph separation
                    if "rules" in header_text.lower() or "circulars" in header_text.lower():
                        value = "\n".join(parts)
                    else:
                        value = " ".join(parts)
                    
                    use_bullets = "gist" in header_text.lower()
                    precomputed = precomputed_bullets_map.get(value) if use_bullets else None
                    
                    # Define style overrides for specific columns
                    style_overrides = None
                    if "rules" in header_text.lower() or "circulars" in header_text.lower():
                        # Custom styling for "Rules / circulars / Notifications / Order" column
                        style_overrides = {
                            "margin_top": 20,      # Keep 20pt top margin
                            "spacing_after": 24,   # Keep 24pt spacing
                        }
                    
                    set_cell_text(cell, value, use_bullets=use_bullets, precomputed_bullets=precomputed, style_overrides=style_overrides)
        
        row_idx += 1

    # Reposition shapes below table (Footer / News Table)
    # Calculate new table bottom
    estimated_table_bottom = shape.top + Inches(current_height)
    
    # We want to position the footer block starting at table_bottom + 0.5 inches
    # But ONLY if this is the last slide (i.e. no remaining rows)
    # If there ARE remaining rows, we might want to push them down anyway to avoid overlap, 
    # or they will be removed in the next step (if we are recursing).
    # Actually, if we are recursing, we remove the News Table from THIS slide.
    # So we only care about positioning if this IS the last slide.
    
    if len(df_remaining) == 0 and footer_shapes:
        logging.info("  Checking footer placement...")
        
        # Calculate where the footer would end up
        estimated_table_bottom = shape.top + Inches(current_height)
        target_footer_top = estimated_table_bottom + Inches(0.5)
        projected_footer_bottom = target_footer_top + Inches(footer_height_inches)
        
        # Check if footer overflows (using 0.3 inch safety margin)
        if projected_footer_bottom > prs.slide_height.inches - 0.3:
            logging.info(f"  Footer overflows (Bottom={projected_footer_bottom / 914400:.2f} > Limit={prs.slide_height.inches - 0.3:.2f}). Moving to NEW slide.")
            
            # 1. Duplicate the slide (contains everything: table + footer)
            new_slide = duplicate_slide(prs, slide_idx)
            new_slide_idx = prs.slides.index(new_slide)
            
            # 2. Remove footer from CURRENT slide
            for s in footer_shapes:
                sp = s.element
                sp.getparent().remove(sp)
                
            # 4. On NEW slide:
            # - Remove the main table completely (since this slide is ONLY for the footer)
            # - Move the footer to the top (where the main table used to be)
            
            new_table_shape = None
            for s in new_slide.shapes:
                if s.has_table and s.top == shape.top and s.left == shape.left:
                    new_table_shape = s
                    break
            
            if new_table_shape:
                # Calculate where the main table starts (to move footer there)
                target_top = new_table_shape.top
                
                # Remove the main table
                sp = new_table_shape.element
                sp.getparent().remove(sp)
                logging.info("  Removed main table from new overflow slide (footer-only slide).")
                
                # 5. Position footer on NEW slide at the top
                
                new_footer_shapes = []
                # We can't rely on shape index or reference after deletion, so re-scan
                # The footer shapes are the ones that were BELOW the table on the original slide.
                # On the new slide, they are still at the bottom.
                # We need to find them and move them up.
                
                # Finding footer shapes on new slide:
                # They should be roughly at the same vertical position as on the original slide
                # But since we just deleted the table, we can just find all shapes that are NOT the title
                # and are below the title.
                # Better: usage of `footer_shapes` list from original slide to identify them by index? 
                # No, indices change.
                # HEURISTIC: Find shapes that look like the footer shapes from the previous slide.
                
                # Actually, we can just look for shapes that are below the title.
                # But let's be safer: we know the footer shapes were below the table bottom on the original slide.
                # On the new slide, they are in the same position.
                
                # Get all shapes on new slide
                all_new_shapes = list(new_slide.shapes)
                
                # Filter for shapes that are likely part of the footer (below where the table WAS)
                # target_top is where the table starts.
                # The footer is currently at `current_footer_top` (approx).
                
                # Let's collect all shapes that are vertically below the target_top
                candidates = []
                for s in all_new_shapes:
                    # Skip Title (usually at top)
                    if s == new_slide.shapes.title:
                        continue
                        
                    # If it's below the target_top, it's a candidate
                    if s.top > target_top:
                        candidates.append(s)
                
                if candidates:
                    # Move them up
                    # Find current top of the block
                    current_block_top = min(s.top for s in candidates)
                    shift = target_top - current_block_top
                    
                    for s in candidates:
                        s.top += shift
                    
                    logging.info(f"  Moved {len(candidates)} footer shapes to top of new slide (Top={target_top/914400:.2f}).")
            
            slide_idx = new_slide_idx
            
        else:
            logging.info("  Footer fits. Repositioning on current slide.")
            try:
                current_footer_top = min(s.top for s in footer_shapes)
                shift = target_footer_top - current_footer_top
                for s in footer_shapes:
                    s.top += shift
                logging.info(f"  Moved footer shapes by {shift.inches if hasattr(shift, 'inches') else shift/914400:.2f} inches")
            except Exception as e:
                logging.warning(f"  Failed to reposition footer: {e}")
    elif len(df_remaining) > 0:
        # If not the last slide, we still need to push them down if they overlap, 
        # just in case they aren't removed (though they should be).
        # Or better, just leave them, they will be removed.
        pass

    # --- RECURSION ---
    total_updates = len(df_current)
    
    if len(df_remaining) > 0:
        logging.info(f"  Creating overflow slide for {len(df_remaining)} remaining rows...")
        
        # Duplicate slide
        new_slide = duplicate_slide(prs, slide_idx)
        new_slide_idx = prs.slides.index(new_slide)
        
        # --- FOOTER REMOVAL LOGIC ---
        # Since we are recursing, THIS slide (slide_idx) is NOT the last one.
        # We should remove ALL footer shapes (including the News Table, backgrounds, etc.) from THIS slide.
        footer_removed = False
        for s in footer_shapes:
            try:
                sp = s.element
                if sp.getparent() is not None:
                    sp.getparent().remove(sp)
                    footer_removed = True
            except Exception as e:
                logging.warning(f"  ⚠ Failed to remove footer shape: {e}")
        
        if footer_removed:
            logging.info(f"  Removed footer elements from non-last Slide {slide_idx+1}")
        # ------------------------
        
        # Find table on new slide
        new_shape = None
        for s in new_slide.shapes:
            # Heuristic: same position as original table
            if abs(s.top - shape.top) < Inches(0.5) and abs(s.left - shape.left) < Inches(0.5):
                if s.has_table:
                    new_shape = s
                    break
        
        if new_shape:
            # Clear data rows from new table (keep header + 1 data row)
            tbl = new_shape.table
            # We need at least 2 rows (Header + 1 Data) for add_table_rows to work
            if len(tbl.rows) > 1:
                rows_to_remove = len(tbl.rows) - 2
                tbl_elem = tbl._tbl
                for _ in range(rows_to_remove):
                    if len(tbl_elem.tr_lst) > 2: # Keep header + 1 data row
                        tbl_elem.remove(tbl_elem.tr_lst[-1])
            
            # Recurse
            # Find the index of the new shape
            new_shape_idx = new_slide.shapes.index(new_shape)
            
            total_updates += paginate_and_fill(
                prs, new_slide_idx, new_shape_idx, 
                df_remaining, headers_map, header_texts, header_to_colspec, precomputed_bullets_map, layout_engine,
                full_df=full_df, # Pass full_df for subtitle updates on new slide
                bottom_margin=bottom_margin  # Preserve margin for recursive overflow slides
            )
        else:
            logging.error("  ✗ Could not find table on duplicated slide")
            
    return total_updates


def fill_using_mapping(prs: Presentation, df: pd.DataFrame, mapping: dict) -> int:
    total_updates = 0
    logging.info(f"DEBUG: fill_using_mapping called with {len(df)} rows")
    
    # Group mappings by (slide_index, shape_index)
    target_groups = {}
    
    for table_spec in mapping.get("tables", []):
        slide_idx = table_spec["slide_index"]
        shape_idx = table_spec["shape_index"]
        key = (slide_idx, shape_idx)
        
        if key not in target_groups:
            target_groups[key] = []
        target_groups[key].append(table_spec)

    # Sort groups by slide index to handle offsets correctly
    sorted_groups = sorted(target_groups.items(), key=lambda x: x[0][0])
    
    slide_offset = 0
    
    # Process each target table
    for (slide_idx, shape_idx), specs in sorted_groups:
        logging.info("=" * 70)
        # Adjust slide index by current offset
        current_slide_idx = slide_idx + slide_offset
        logging.info(f"Processing Target: Original Slide {slide_idx + 1} -> Current Slide {current_slide_idx + 1}, Shape {shape_idx}")
        logging.info("=" * 70)
        
        # Combine data for this target
        combined_data_parts = []
        
        # Get headers map from the first spec (assuming all specs for same table have same headers)
        headers_map = specs[0]["headers"]
        
        for spec in specs:
            source_filter = spec.get("source_filter", {})
            
            # Filter data
            data_subset = df.copy()
            for col, val in source_filter.items():
                if col in data_subset.columns:
                    if isinstance(val, list):
                        data_subset = data_subset[data_subset[col].isin(val)]
                    else:
                        data_subset = data_subset[data_subset[col] == val]
            
            if len(data_subset) > 0:
                combined_data_parts.append(data_subset)
                logging.info(f"  Added {len(data_subset)} rows for {source_filter}")

        if not combined_data_parts:
            logging.warning("  No data found for this target.")
            continue
            
        combined_df = pd.concat(combined_data_parts, ignore_index=True)
        
        # Extract header texts from the actual table to map correctly
        if current_slide_idx < len(prs.slides):
            slide = prs.slides[current_slide_idx]
            
            shape = None
            # Try specific index first
            if shape_idx < len(slide.shapes):
                s = slide.shapes[shape_idx]
                if s.has_table:
                    shape = s
            
            # Fallback: Search for any table
            if not shape:
                logging.warning(f"  ⚠ Shape {shape_idx} not found or not a table. Searching for first table...")
                for i, s in enumerate(slide.shapes):
                    if s.has_table:
                        shape = s
                        shape_idx = i # Update index for paginate_and_fill
                        logging.info(f"  ✓ Found table at index {i}")
                        break
            
            if shape and shape.has_table:
                header_row = shape.table.rows[0]
                header_texts = [cell.text.strip() for cell in header_row.cells]
                
                # Map headers
                header_to_colspec = {}
                for header_text in header_texts:
                    normalized = " ".join(header_text.split())
                    colspec = headers_map.get(header_text) or headers_map.get(normalized)
                    if colspec:
                        header_to_colspec[header_text] = colspec
                
                # Precompute bullets
                precomputed_bullets_map = {}
                # Identify which columns require bullets
                bullet_cols = []
                for header_text in header_texts:
                    if "gist" in header_text.lower():
                        colspec = header_to_colspec.get(header_text)
                        if colspec and isinstance(colspec, list):
                            bullet_cols.extend(colspec)
                
                if bullet_cols:
                    from src.services.ppt.engine import convert_to_bullet_points, format_date
                    import concurrent.futures
                    texts_to_bulletize = set()
                    
                    for _, row in combined_df.iterrows():
                        is_section_header = row.get('IsSectionHeader', False)
                        if is_section_header:
                            continue
                        # Build the text that will be bulletized
                        parts = []
                        for col_name in bullet_cols:
                            if col_name in combined_df.columns:
                                v = row[col_name]
                                if not pd.isna(v):
                                    if col_name == "IssueDate":
                                        parts.append(format_date(v))
                                    else:
                                        parts.append(str(v))
                        val = " ".join(parts)
                        if val.strip():
                            texts_to_bulletize.add(val)
                    
                    if texts_to_bulletize:
                        logging.info(f"  Precomputing bullets for {len(texts_to_bulletize)} items concurrently...")
                        with concurrent.futures.ThreadPoolExecutor(max_workers=5) as executor:
                            future_to_text = {
                                executor.submit(convert_to_bullet_points, text): text 
                                for text in texts_to_bulletize
                            }
                            for future in concurrent.futures.as_completed(future_to_text):
                                text = future_to_text[future]
                                try:
                                    bullets = future.result()
                                    precomputed_bullets_map[text] = bullets
                                except Exception as exc:
                                    logging.error(f"  ⚠ Bullet generation generated an exception: {exc}")
                        logging.info("  ✓ Precomputation complete.")
                
                # Track slides before filling
                slides_before = len(prs.slides)
                
                # Call recursive filler
                # IFSCA slides have very long URLs that the height estimator under-counts.
                # Use a larger bottom_margin so overflow triggers earlier, preventing visual overflow.
                is_ifsca = any(
                    "ifsca" in str(spec.get("source_filter", {}).get("Verticals", "")).lower() or
                    "ifsc" in str(spec.get("source_filter", {}).get("Verticals", "")).lower()
                    for spec in specs
                )
                ifsca_margin = 1.5 if is_ifsca else 0.5
                
                updates = paginate_and_fill(
                    prs, current_slide_idx, shape_idx, combined_df, 
                    headers_map, header_texts, header_to_colspec, precomputed_bullets_map,
                    full_df=combined_df, # Pass full dataframe for subtitle counts
                    bottom_margin=ifsca_margin
                )
                total_updates += updates
                
                # Update offset based on new slides added
                slides_added = len(prs.slides) - slides_before
                if slides_added > 0:
                    slide_offset += slides_added
                    logging.info(f"  Added {slides_added} overflow slides. New offset: {slide_offset}")

    return total_updates


def read_excel_files_from_directory(directory_path: str) -> pd.DataFrame:
    """
    Read all Excel files from a directory and combine them into a single DataFrame.
    Each file is treated as a different vertical based on FILE_VERTICAL_MAPPING.
    Each sheet within a file is treated as a subdomain.
    
    Args:
        directory_path: Path to directory containing Excel files
        
    Returns:
        Combined DataFrame with SourceFile, Subdomain, and Verticals columns
    """
    if not os.path.exists(directory_path):
        logging.error(f"Directory not found: {directory_path}")
        return pd.DataFrame()
    
    if not os.path.isdir(directory_path):
        logging.error(f"Path is not a directory: {directory_path}")
        return pd.DataFrame()
    
    # Find all Excel files (exclude temp files)
    excel_files = glob.glob(os.path.join(directory_path, "*.xlsx"))
    excel_files = [f for f in excel_files if not os.path.basename(f).startswith("~$")]
    
    if not excel_files:
        logging.warning(f"No Excel files found in {directory_path}")
        return pd.DataFrame()
    
    logging.info(f"Found {len(excel_files)} Excel file(s) in {directory_path}")
    
    all_dfs = []
    
    for excel_path in excel_files:
        file_basename = os.path.splitext(os.path.basename(excel_path))[0]
        
        # Map file name to vertical category
        vertical_name = FILE_VERTICAL_MAPPING.get(file_basename, file_basename)
        
        logging.info(f"Processing file: {file_basename} → Vertical: {vertical_name}")
        
        try:
            # Read all sheets from this file
            excel_file = pd.ExcelFile(excel_path)
            sheets = excel_file.sheet_names
            logging.info(f"  Found {len(sheets)} sheet(s): {sheets}")
            
            for sheet_name in sheets:
                sheet_df = pd.read_excel(excel_path, sheet_name=sheet_name)
                
                # Add tracking columns
                sheet_df['SourceFile'] = file_basename
                sheet_df['Subdomain'] = sheet_name
                
                # Add or override Verticals column based on file mapping
                sheet_df['Verticals'] = vertical_name
                
                all_dfs.append(sheet_df)
                logging.info(f"    - {sheet_name}: {len(sheet_df)} rows")
                
        except Exception as e:
            logging.error(f"Error reading {excel_path}: {e}")
            continue
    
    if not all_dfs:
        logging.error("No data could be read from any Excel files")
        return pd.DataFrame()
    
    # Combine all DataFrames
    combined_df = pd.concat(all_dfs, ignore_index=True)
    
    # Ensure column names are stripped
    combined_df.columns = combined_df.columns.str.strip()
    
    logging.info(f"Total combined: {len(combined_df)} rows from {len(excel_files)} file(s)")
    
    return combined_df


def main():
    parser = argparse.ArgumentParser(description="Smart PPT Agent")
    parser.add_argument("--excel", help="Path to a single Excel file (for backward compatibility)", default=None)
    parser.add_argument("--excel-dir", help="Path to directory containing multiple Excel files", default=None)
    parser.add_argument("--ppt", help="Path to the PowerPoint template", default=None)
    parser.add_argument("--local-only", help="Skip MinIO and use only local files", action="store_true")
    parser.add_argument("--log-level", help="Logging level (DEBUG, INFO, WARNING, ERROR)", default="INFO")
    parser.add_argument("--log-file", help="Log file path (optional)", default=None)
    args = parser.parse_args()

    # Setup logging
    log_level = getattr(logging, args.log_level.upper(), logging.INFO)
    setup_logging(log_level=log_level, log_file=args.log_file)

    # Initialize variables
    excel_dir = None
    excel_path = None

    # MinIO Integration: Fetch Excel files from object storage
    # Skip MinIO if --local-only flag is set
    if args.local_only:
        logging.info("--local-only flag set. Skipping MinIO, using local files only.")
        minio_client = None
    else:
        logging.info("Initializing MinIO Client...")
        try:
            minio_client = MinioClient()
            logging.info(f"Connected to MinIO at {minio_client.client._base_url}")
        except Exception as e:
            logging.error(f"Failed to initialize MinIO client: {e}")
            # Continue to fallback logic
            minio_client = None

    # List Excel files from MinIO
    excel_files = []
    if minio_client:
        # Import date utilities
        from src.utils.date_utils import get_current_week_folder_name, log_week_info
        
        # Log current week information
        log_week_info()
        
        # Calculate current week folder name
        current_week_folder = get_current_week_folder_name()
        logging.info(f"Looking for previous week folder (this week's data): {current_week_folder}")
        
        try:
            # First, try to fetch files from the current week folder
            excel_files = minio_client.list_excel_files(prefix=current_week_folder)
            
            if excel_files:
                logging.info(f"✓ Found {len(excel_files)} file(s) in current week folder: {current_week_folder}")
            else:
                logging.warning(f"⚠ Current week folder '{current_week_folder}' exists but contains no Excel files.")
                logging.info("Falling back to latest available folder...")
                
                # Fallback: List all folders and find the latest one
                objects = minio_client.client.list_objects(minio_client.bucket_name, prefix="weekly_outputs/", recursive=False)
                
                folders = []
                for obj in objects:
                    if obj.object_name.endswith('/'):
                        folders.append(obj.object_name)
                
                if folders:
                    # Sort folders to find the latest one
                    latest_folder = sorted(folders)[-1]
                    logging.info(f"Using latest available folder: {latest_folder}")
                    excel_files = minio_client.list_excel_files(prefix=latest_folder)
                else:
                    logging.warning("No weekly folders found. Scanning entire bucket...")
                    excel_files = minio_client.list_excel_files()
                    
        except Exception as e:
            logging.error(f"Error accessing current week folder: {e}")
            logging.info("Falling back to full bucket scan...")
            excel_files = minio_client.list_excel_files()
    
    if not excel_files:
        logging.warning("No Excel files found in MinIO bucket (or MinIO unavailable).")
        # Continue to fallback logic

    # Process each Excel file
    all_data_frames = []
    
    for file_obj_name in excel_files:
        logging.info(f"Processing file from MinIO: {file_obj_name}")
        
        try:
            # Download file content
            file_content = minio_client.get_file_content(file_obj_name)
            
            # Read Excel from memory
            # Use ExcelFile to read all sheets (subdomains)
            excel_file = pd.ExcelFile(file_content)
            sheet_names = excel_file.sheet_names
            logging.info(f"  Found {len(sheet_names)} sheet(s): {sheet_names}")
            
            file_dfs = []
            for sheet_name in sheet_names:
                sheet_df = pd.read_excel(excel_file, sheet_name=sheet_name)
                # Add subdomain column
                sheet_df['Subdomain'] = sheet_name
                file_dfs.append(sheet_df)
                logging.info(f"    - Sheet '{sheet_name}': {len(sheet_df)} rows")
            
            # Combine sheets for this file
            if file_dfs:
                df = pd.concat(file_dfs, ignore_index=True)
            else:
                continue
            
            # Determine vertical from filename (basename without extension)
            file_basename = os.path.basename(file_obj_name)
            filename_no_ext = os.path.splitext(file_basename)[0]
            
            vertical = FILE_VERTICAL_MAPPING.get(filename_no_ext, filename_no_ext)
            logging.info(f"  Mapped '{file_basename}' to Vertical: '{vertical}'")
            
            # Add or override Verticals column based on file mapping
            # We use 'Verticals' (plural) to match the rest of the system and overwrite any existing value
            df["Verticals"] = vertical
                
            all_data_frames.append(df)
            
        except Exception as e:
            logging.error(f"Failed to process file '{file_obj_name}': {e}")
            continue
            
    if not all_data_frames:
        logging.warning("No valid data found in MinIO. Falling back to local files...")
        
        # Determine input mode: multi-file directory or single file
        excel_dir = args.excel_dir
        excel_path = args.excel
        
        # If neither provided, try to auto-detect excels/ directory
        if not excel_dir and not excel_path:
            if os.path.exists("excels") and os.path.isdir("excels"):
                excel_dir = "excels"
                logging.info(f"Auto-detected Excel directory: {excel_dir}")
            else:
                # Fall back to single-file detection
                local_excel_files = glob.glob("*.xlsx")
                if local_excel_files:
                    excel_path = local_excel_files[0]
                    logging.info(f"Auto-detected Excel file: {excel_path}")
                else:
                    logging.error("Error: No Excel file or directory found (Local or MinIO).")
                    return

        # Read Excel data based on mode
        if excel_dir:
            # MULTI-FILE MODE: Read all Excel files from directory
            logging.info(f"Mode: Multi-file directory ({excel_dir})")
            df_local = read_excel_files_from_directory(excel_dir)
            
            if df_local.empty:
                logging.error("No data could be read from directory")
                return
            
            all_sheets = df_local['Subdomain'].unique().tolist() if 'Subdomain' in df_local.columns else []
            all_data_frames.append(df_local)
                    
        elif excel_path:
            # SINGLE-FILE MODE
            logging.info(f"Mode: Single file ({excel_path})")
            if not os.path.exists(excel_path):
                logging.error(f"Error: Excel file not found at {excel_path}")
                return
                
            try:
                excel_file = pd.ExcelFile(excel_path)
                sheet_names = excel_file.sheet_names
                logging.info(f"Found {len(sheet_names)} sheet(s) in Excel: {sheet_names}")
                
                file_dfs = []
                for sheet_name in sheet_names:
                    sheet_df = pd.read_excel(excel_file, sheet_name=sheet_name)
                    sheet_df['Subdomain'] = sheet_name
                    file_dfs.append(sheet_df)
                    logging.info(f"  - {sheet_name}: {len(sheet_df)} rows")
                
                if file_dfs:
                    df = pd.concat(file_dfs, ignore_index=True)
                    # For single file, assume default vertical or derived from filename
                    filename = os.path.basename(excel_path)
                    filename_no_ext = os.path.splitext(filename)[0]
                    vertical = FILE_VERTICAL_MAPPING.get(filename_no_ext, filename_no_ext)
                    
                    # Use 'Verticals' (plural) to match the rest of the system
                    if "Verticals" not in df.columns:
                        df["Verticals"] = vertical
                    all_data_frames.append(df)
                    
            except Exception as e:
                logging.error(f"Failed to read {excel_path}: {e}")
                return

    if not all_data_frames:
        logging.error("No valid data found in Excel files (MinIO or Local).")
        return

    # Combine all dataframes
    df = pd.concat(all_data_frames, ignore_index=True)
    logging.info(f"Combined Data: {len(df)} rows from {len(all_data_frames)} files")

    # Filter to only include allowed verticals (SEBI, AIF, Listed Companies, RBI)
    ALLOWED_VERTICALS = ["SEBI", "Alternate Investment Funds", "Others during the week", "IFSCA", "RBI", "Insolvency and Bankruptcy Code", "Companies Act"]
    
    if "Verticals" in df.columns:
        rows_before = len(df)
        df = df[df["Verticals"].isin(ALLOWED_VERTICALS)]
        rows_after = len(df)
        rows_filtered = rows_before - rows_after
        
        logging.info("=" * 70)
        logging.info("DOMAIN FILTERING")
        logging.info("=" * 70)
        logging.info(f"Allowed Verticals: {ALLOWED_VERTICALS}")
        logging.info(f"Rows before filtering: {rows_before}")
        logging.info(f"Rows after filtering: {rows_after}")
        logging.info(f"Rows excluded: {rows_filtered}")
        logging.info("=" * 70)
        
        if rows_after == 0:
            logging.error("No data remaining after filtering. Check if Verticals column values match allowed list.")
            return
    else:
        logging.warning("'Verticals' column not found. Skipping domain filtering.")
        
    # Filter out specific subdomains for RBI vertical
    RBI_EXCLUDED_SUBDOMAINS = ["Rules","Regulations"]


    
    if "Verticals" in df.columns and "Subdomain" in df.columns:
        rows_before = len(df)
        
        # Create mask for exclusion
        rbi_mask = (df["Verticals"] == "RBI") & (df["Subdomain"].isin(RBI_EXCLUDED_SUBDOMAINS))
        
        df = df[~rbi_mask]
        rows_after = len(df)
        rows_excluded = rows_before - rows_after
        
        if rows_excluded > 0:
            logging.info("=" * 70)
            logging.info("RBI SUBDOMAIN FILTERING")
            logging.info("=" * 70)
            logging.info(f"Excluded Subdomains for RBI: {RBI_EXCLUDED_SUBDOMAINS}")
            logging.info(f"Rows before filtering: {rows_before}")
            logging.info(f"Rows after filtering: {rows_after}")
            logging.info(f"RBI rows excluded: {rows_excluded}")
            logging.info("=" * 70)


    ppt_path = args.ppt
    if not ppt_path:
        ppt_files = glob.glob("*.pptx")
        # Filter out temporary files or output files if possible, but for now just take the first one that isn't a temp file
        ppt_files = [f for f in ppt_files if not f.startswith("~$") and "filled" not in f]
        if ppt_files:
            ppt_path = ppt_files[0]
            logging.info(f"Auto-detected PPT template: {ppt_path}")
        else:
            logging.error("Error: No PPT template found. Please provide one using --ppt or place one in the current directory.")
            return
            
    if not os.path.exists(ppt_path):
        logging.error(f"Error: PPT file not found at {ppt_path}")
        return
    
    logging.info("=" * 70)
    logging.info("SMART PPT AGENT - MINIO INTEGRATION + MANUAL SLIDE MAPPING")
    logging.info("=" * 70)
    
    # (Skip local file reading block)
    # Data is already loaded from MinIO
    if df.empty:
        logging.error("No data loaded from MinIO")
        return
        
    all_sheets = df['Subdomain'].unique().tolist() if 'Subdomain' in df.columns else []


    

    # Ensure all column names are stripped of whitespace
    df.columns = df.columns.str.strip()

    logging.info(f"Excel: {len(df)} total rows")
    
    # Show source files breakdown (for multi-file mode)
    if 'SourceFile' in df.columns:
        source_files = df['SourceFile'].unique()
        logging.info(f"Source files: {list(source_files)}")
        for source_file in source_files:
            count = len(df[df['SourceFile'] == source_file])
            logging.info(f"  - {source_file}: {count} rows")
    
    # Show verticals breakdown
    if 'Verticals' in df.columns:
        verticals = df['Verticals'].unique()
        logging.info(f"Verticals found: {list(verticals)}")
        for vertical in verticals:
            count = len(df[df['Verticals'] == vertical])
            logging.info(f"  - {vertical}: {count} rows")
    
    # Show subdomains breakdown
    if 'Subdomain' in df.columns:
        subdomains = df['Subdomain'].unique()
        logging.info(f"Subdomains found: {list(subdomains)}")
        for subdomain in subdomains:
            count = len(df[df['Subdomain'] == subdomain])
            logging.info(f"  - {subdomain}: {count} rows")

    
    # Load presentation
    prs = Presentation(ppt_path)
    
    # Extract template structure
    template_spec = extract_template_spec(ppt_path, max_slides=20)
    
    # Build data profile
    # Build data profile directly from loaded DataFrame
    # This works for MinIO, multi-file, and single-file modes since df is always populated
    data_profile = {
        "columns": df.columns.tolist(),
        "unique_verticals": df["Verticals"].unique().tolist() if "Verticals" in df.columns else [],
        "unique_subdomains": df["Subdomain"].unique().tolist() if "Subdomain" in df.columns else [],
        "unique_source_files": df["SourceFile"].unique().tolist() if "SourceFile" in df.columns else [],
        "num_rows": len(df)
    }
    
    # Get LLM
    llm = get_llm()
    logging.info(f"DEBUG: len(df) before mapping: {len(df)}")
    
    # Get mapping (AI or fallback)
    mapping = ask_llm_for_mapping(llm, template_spec, data_profile) if llm else get_fallback_mapping(template_spec, data_profile)
    
    logging.info("--- Generated Mapping ---")
    logging.debug(json.dumps(mapping, indent=2))
    
    logging.info(f"DEBUG: len(df) before fill: {len(df)}")
    
    # Fill using mapping
    updates = fill_using_mapping(prs, df, mapping)
    
    # Save output
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    output_path = f"tejomaya_filled_{timestamp}.pptx"
    prs.save(output_path)
    
    logging.info("=" * 70)
    logging.info(f"✓ SUCCESS! Filled {updates} total rows across all sectors")
    logging.info(f"✓ Saved: {output_path}")
    logging.info("=" * 70)
    logging.info(f"open {output_path}")

if __name__ == "__main__":
    main()
