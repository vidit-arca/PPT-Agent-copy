import json
from pptx import Presentation

def extract_template_spec(ppt_path: str, max_slides: int = 20) -> dict:
    prs = Presentation(ppt_path)
    
    template_spec = {
        "total_slides": len(prs.slides),
        "slides": []
    }
    
    # IMPORTANT: do NOT slice prs.slides; iterate with index and break
    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx >= max_slides:
            break

        slide_info = {
            "slide_index": slide_idx,
            "shapes": []
        }
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_info = {
                "shape_index": shape_idx,
                "name": getattr(shape, "name", f"Shape_{shape_idx}"),
                "type": None,
                "text": None,
                "is_table": False,
                "table_headers": []
            }
            
            # Text box
            if hasattr(shape, "text_frame") and shape.text_frame:
                txt = shape.text_frame.text.strip()
                if txt:
                    shape_info["type"] = "text_box"
                    shape_info["text"] = txt
            
            # Table
            if shape.has_table:
                shape_info["type"] = "table"
                shape_info["is_table"] = True
                table = shape.table
                headers = []
                if len(table.rows) > 0:
                    for cell in table.rows[0].cells:
                        headers.append(cell.text.strip())
                shape_info["table_headers"] = headers
            
            # Only keep shapes that have meaningful info
            if shape_info["type"] is not None:
                slide_info["shapes"].append(shape_info)
        
        template_spec["slides"].append(slide_info)
    
    return template_spec

if __name__ == "__main__":
    import glob
    import os
    
    ppt_files = glob.glob("*.pptx")
    # Filter out temporary files
    ppt_files = [f for f in ppt_files if not f.startswith("~$")]
    
    if ppt_files:
        ppt_path = ppt_files[0]
        print(f"Using PPT: {ppt_path}")
        spec = extract_template_spec(ppt_path, max_slides=20)
        print(json.dumps(spec, indent=2))
    else:
        print("No PPT file found in current directory.")
