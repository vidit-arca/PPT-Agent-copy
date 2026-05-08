import logging
import re
import pandas as pd
from copy import deepcopy
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from lxml import etree

from src.services.text.processor import AITextProcessor

def duplicate_slide(prs, source_slide_index):
    """
    Duplicate the slide at source_slide_index.
    Inserts the new slide right after the source slide.
    Returns the new slide object.
    """
    source_slide = prs.slides[source_slide_index]
    slide_layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    # Copy shapes
    copied_count = 0
    for shape in source_slide.shapes:
        # Create a copy of the element
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
        copied_count += 1
        
        # Copy relationships (images, hyperlinks, etc.)
        # We need to check if the shape has any relationships and copy them to the new slide's part
        # This is a bit low-level but necessary for images to show up
        
        # Iterate over all relationships in the source slide part
        # and check if this shape references any of them.
        # A shape references a relationship via r:embed or r:link attributes in its XML.
        
        # Simple approach: Copy ALL relationships from source slide to new slide
        # This might result in unused relationships but ensures nothing is missing.
        # However, we need to map old rIds to new rIds if they change.
        # But if we just add them, python-pptx assigns new rIds.
        # So we must update the shape's XML to point to the new rId.
        
        # Let's scan the shape's XML for r:embed or r:link
        for elem in new_el.iter():
            for attr, value in elem.attrib.items():
                if "embed" in attr or "link" in attr:
                    # This is a relationship ID (e.g., "rId2")
                    # Find the relationship in the source slide
                    try:
                        rel = source_slide.part.rels[value]
                        target = rel.target_part
                        
                        # Add this relationship to the new slide
                        # We need to determine the relationship type
                        rel_type = rel.reltype
                        
                        # Add relationship to new slide part
                        new_rid = new_slide.part.rels.get_or_add(rel_type, target)
                        
                        # Update the attribute in the new shape to use the new rId
                        elem.attrib[attr] = new_rid
                    except KeyError:
                        # Relationship not found (maybe internal or standard?)
                        pass
    
    # Move the new slide to the position right after the source slide
    # The new slide was added at the end, we need to move it
    # We work with the slide ID list in the presentation
    xml_slides = prs.slides._sldIdLst
    
    # Find the slide IDs
    slides = list(prs.slides)
    new_slide_idx = slides.index(new_slide)
    
    # Only move if it's not already in the right position
    if new_slide_idx != source_slide_index + 1:
        # Get the slide ID element for the new slide
        slide_id_elem = xml_slides[new_slide_idx]
        
        # Remove it from current position
        xml_slides.remove(slide_id_elem)
        
        # Insert at target position (after source slide)
        xml_slides.insert(source_slide_index + 1, slide_id_elem)
    
    # Remove any empty placeholder shapes that might have been added by the layout
    # These can cause visual artifacts and overflow issues
    placeholders_to_remove = []
    for shape in new_slide.shapes:
        if shape.is_placeholder:
            # Check if it's empty placeholder (Title or Content placeholders that are blank)
            if hasattr(shape, 'text_frame') and shape.text_frame:
                if not shape.text_frame.text.strip():
                    # Empty text placeholder - remove it
                    placeholders_to_remove.append(shape)
    
    for placeholder in placeholders_to_remove:
        try:
            remove_shape(new_slide, placeholder)
            logging.debug(f"  Removed empty placeholder from duplicate slide")
        except Exception as e:
            pass
    
    return new_slide


def get_shapes_below(slide, threshold_top):
    """Return a list of shapes positioned below the threshold_top."""
    shapes_below = []
    for shape in slide.shapes:
        if shape.top > threshold_top:
            shapes_below.append(shape)
    return shapes_below


def remove_shape(slide, shape):
    """Remove a shape from a slide."""
    sp = shape.element
    sp.getparent().remove(sp)


def add_table_rows(table, num_rows_to_add: int):
    """Add rows by duplicating the last row."""
    if len(table.rows) < 2 or num_rows_to_add <= 0:
        return

    last_row_index = len(table.rows) - 1
    last_row = table.rows[last_row_index]
    tbl = table._tbl

    for _ in range(num_rows_to_add):
        new_tr = deepcopy(last_row._tr)
        for tc in new_tr.tc_lst:
            for p in tc.txBody.p_lst:
                for r in p.r_lst:
                    r.t.text = ""
        tbl.append(new_tr)


def remove_table_rows(table, num_rows_to_remove: int):
    """Remove rows from the end of the table."""
    if num_rows_to_remove <= 0:
        return

    tbl = table._tbl
    # Ensure we don't remove the header row (keep at least 1 row)
    current_rows = len(table.rows)
    # We can remove up to current_rows - 1 (keep header)
    max_removable = current_rows - 1
    
    to_remove = min(num_rows_to_remove, max_removable)
    
    if to_remove <= 0:
        return
        
    for _ in range(to_remove):
        # Remove last row
        if len(tbl.tr_lst) > 1:
            tbl.remove(tbl.tr_lst[-1])


def move_shapes_below(slide, table_shape, added_rows_count):
    """Move shapes below the table down to prevent overlap."""
    if added_rows_count <= 0:
        return

    offset_down = Inches(0.3 * added_rows_count)
    table_bottom = table_shape.top + table_shape.height

    logging.debug(f"  Moving shapes below table down by {offset_down.inches:.2f} inches...")

    moved_count = 0
    for shape in slide.shapes:
        if shape == table_shape:
            continue
        if shape.top > table_bottom:
            shape.top += offset_down
            moved_count += 1

    if moved_count > 0:
        logging.info(f"  ✓ Moved {moved_count} shape(s) down")


def format_date(value):
    """Format date to DD-MM-YY if it's a datetime object or string."""
    if isinstance(value, pd.Timestamp):
        return value.strftime("%d-%m-%y")
    try:
        # Try parsing string date
        dt = pd.to_datetime(value)
        return dt.strftime("%d-%m-%y")
    except:
        return str(value)


def convert_to_bullet_points(text: str) -> list:
    """Convert text to bullet points using AI (Phase 2) with rule-based fallback.
    
    Args:
        text: The text to convert
        
    Returns:
        List of bullet point strings
    """
    # Try AI-powered bullet generation first
    try:
        processor = AITextProcessor()
        response = processor.generate_bullets(text, max_bullets=5)
        
        # Use AI bullets if quality is good
        if response.summary_quality >= 7:
            logging.info(f"✓ AI bullets: {len(response.bullets)} points, quality {response.summary_quality}/10")
            return response.bullets
        else:
            logging.warning(f"⚠ AI quality low ({response.summary_quality}/10), using fallback")
            return _fallback_bullet_split(text)
            
    except Exception as e:
        logging.warning(f"⚠ AI bullet generation failed: {e}, using fallback")
        return _fallback_bullet_split(text)


def _fallback_bullet_split(text: str) -> list:
    """Rule-based bullet splitting as fallback (original Phase 1 logic)."""
    # Remove surrounding quotes first
    text = text.strip().strip('"\'')
    text = ' '.join(text.split())  # Clean up extra whitespace
    
    # Try different splitting strategies
    if ';' in text:
        points = [p.strip() for p in text.split(';') if p.strip()]
    elif '. ' in text:
        points = re.split(r'\.\s+(?=[A-Z])', text)
        points = [p.strip() for p in points if p.strip()]
        points = [p if p.endswith('.') else p + '.' for p in points[:-1]] + [points[-1]] if points else []
    elif '\n' in text:
        points = [p.strip() for p in text.split('\n') if p.strip()]
    else:
        points = [text]
    
    # Clean up points
    cleaned_points = []
    for point in points:
        point = point.strip()
        point = re.sub(r'^[•\-\*]\s*', '', point)
        point = point.strip('"\'')
        point = ' '.join(point.split())
        if point:
            cleaned_points.append(point)
    
    return cleaned_points


from pptx.enum.text import MSO_ANCHOR

def set_cell_text(cell, value: str, use_bullets: bool = False, precomputed_bullets: list = None, style_overrides: dict = None):
    """Write text to cell with AI-powered formatting (Phase 2).
    
    Args:
        cell: The table cell to write to
        value: The text value to write
        use_bullets: If True, convert text to bullet points with AI formatting
        precomputed_bullets: Optional list of bullets to use directly (skips generation)
        style_overrides: Optional dict of style overrides (font_size, spacing_before, spacing_after, margin_top, etc.)
    """
    if use_bullets:
        # 1. Use precomputed bullets or generate new ones
        if precomputed_bullets:
            bullet_points = precomputed_bullets
        else:
            bullet_points = convert_to_bullet_points(value)
        
        if not bullet_points:
            cell.text = ""
            return
        
    # Default formatting
    style = {
        "font_size": 12,
        "spacing_before": 3,
        "spacing_after": 10,
        "margin_left": 5,
        "margin_right": 5,
        "margin_top": 20,
        "margin_bottom": 3
    }
    
    # Apply overrides if provided
    if style_overrides:
        style.update(style_overrides)
        logging.debug(f"Applying style overrides: {style_overrides}")

    # Apply consistent margins to CELL directly (not just text_frame)
    cell.margin_left = Pt(style["margin_left"])
    cell.margin_right = Pt(style["margin_right"])
    cell.margin_top = Pt(style["margin_top"])
    cell.margin_bottom = Pt(style["margin_bottom"])
    cell.vertical_anchor = MSO_ANCHOR.TOP

    # Also set text_frame margins just in case, but cell margins should govern
    text_frame = cell.text_frame
    text_frame.margin_left = Pt(0)
    text_frame.margin_right = Pt(0)
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)

    if use_bullets:
        # 1. Use precomputed bullets or generate new ones
        if precomputed_bullets:
            bullet_points = precomputed_bullets
        else:
            bullet_points = convert_to_bullet_points(value)
        
        if not bullet_points:
            cell.text = ""
            return
        
        logging.debug(f"Using fixed formatting: {style['font_size']}pt, spacing {style['spacing_before']}/{style['spacing_after']}pt")
        
        # 3. Apply consistent formatting
        text_frame.clear()
        
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = point
            p.level = 0
            p.space_before = Pt(style["spacing_before"])
            p.space_after = Pt(style["spacing_after"])
            
            # Apply hanging indent using XML directly since paragraph_format might be missing
            # 18pt = 228600 EMUs
            # -12pt = -152400 EMUs
            pPr = p._element.get_or_add_pPr()
            pPr.set('marL', '228600')
            pPr.set('indent', '-152400')
            
            # Apply bullet formatting using XML
            pPr = p._element.get_or_add_pPr()
            
            for child in list(pPr):
                if child.tag.endswith('buNone') or child.tag.endswith('buAutoNum'):
                    pPr.remove(child)
            
            from lxml import etree
            buChar = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
            buChar.set('char', '•')
            
            buFont = etree.SubElement(pPr, '{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
            buFont.set('typeface', 'Arial')
            
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(style["font_size"])
                run.font.name = "Arial"
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    else:
        # Non-bullet text - also use consistent 10pt font and spacing
        # Clear existing paragraphs to ensure clean state? No, just use the first one.
        text_frame.clear() # Clear to ensure we start fresh for multi-paragraph support
        
        # Split by newline to support multi-paragraph content (e.g. Rules column)
        paragraphs = value.split('\n')
        
        for i, para_text in enumerate(paragraphs):
            if not para_text.strip() and len(paragraphs) > 1:
                continue # Skip empty lines if we have multiple
                
            if i == 0:
                # text_frame.clear() might leave one empty paragraph or remove all?
                # python-pptx text_frame.clear() removes all text but leaves one empty paragraph usually?
                # Actually clear() removes all paragraphs. So we need to add one.
                # Wait, let's check python-pptx behavior. 
                # If we use clear(), we need to add_paragraph()
                p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
            else:
                p = text_frame.add_paragraph()
                
            p.text = para_text
            p.space_before = Pt(style["spacing_before"])
            p.space_after = Pt(style["spacing_after"])
            
            # Apply consistent 10pt font to all cells
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(style["font_size"])  # Consistent 10pt
                run.font.name = "Arial"
                run.font.color.rgb = RGBColor(0, 0, 0)


def update_slide_subtitle(slide, df):
    """Update the subtitle on the slide with counts of each subdomain."""
    
    # Detect vertical type
    is_others_vertical = False
    is_rbi_vertical = False
    is_companies_act_vertical = False
    is_ibbi_vertical = False
    if "Verticals" in df.columns:
        verticals = df["Verticals"].unique()
        is_others_vertical = "Others during the week" in verticals
        is_rbi_vertical = "RBI" in verticals
        is_companies_act_vertical = "Companies Act" in verticals
        is_ibbi_vertical = "Insolvency and Bankruptcy Code" in verticals
    
    # Define categories based on vertical type
    if is_others_vertical:
        # Categories for "Others during the week" slide
        categories = {
            "ICAI": ["ICAI"],
            "DGFT": ["DGFT"],
            "NCLT": ["NCLT"],
            "ICSI": ["ICSI"],
            "NSE": ["NSE", "Circular-NSE"],
            "BSE": ["BSE", "Circular-BSE"],  # BSE as separate category
            "Others": ["Others"]  # Only actual "Others" items
        }
        category_order = ["ICAI", "DGFT", "NCLT", "ICSI", "NSE", "BSE", "Others"]
    elif is_rbi_vertical:
        # Categories for RBI slide
        categories = {
            "Notifications": ["Notifications", "Notification"],
            "Master Directions": ["Master Directions", "Master Direction"],
            "Circulars": ["Circulars", "Circular"],
            "Press Release": ["Press Release", "Press Releases"]
        }
        category_order = ["Notifications", "Master Directions", "Circulars", "Press Release"]
    elif is_companies_act_vertical:
        # Categories for "Companies Act 2013 - MCA" slide
        categories = {
            "Notifications": ["Notifications", "Notification"],
            "Rules": ["Rules", "Rule"],
            "Circulars": ["Circulars", "Circular"],
            "Orders": ["Orders", "Order"],
            "Press Release": ["Press Release", "Press Releases"],
            "Public Notice": ["Public Notice", "Public Notices"]
        }
        category_order = ["Notifications", "Rules", "Circulars", "Orders", "Press Release", "Public Notice"]
    elif is_ibbi_vertical:
        # Categories for "Insolvency and Bankruptcy Code" slide
        categories = {
            "Press Release": ["Press Release", "Press Releases"],
            "Master Direction": ["Master Direction", "Master Directions"],
            "Notification": ["Notification", "Notifications"],
            "Circulars": ["Circulars", "Circular"],
            "Discussion Paper": ["Discussion Paper", "Discussion Papers"]
        }
        category_order = ["Press Release", "Master Direction", "Notification", "Circulars", "Discussion Paper"]
    else:
        # Standard categories for other verticals (SEBI, AIF, IFSCA, etc.)
        categories = {
            "Circulars": ["Circulars", "Circular", "Circular-NSE", "Circular-BSE"],
            "Consultation Paper": ["Consultation Paper", "Consulatation Paper"],
            "Regulation": ["Regulation", "Regulations"],
            "Master Circular": ["Master Circular", "Master circular"],
            "Press Release": ["Press Release"],
            "Adjudication Orders": ["Adjudication Orders", "Adjudication Order", "Orders", "Order"]
        }
        category_order = ["Circulars", "Consultation Paper", "Regulation", "Master Circular", "Press Release", "Adjudication Orders"]
    
    # Calculate counts
    counts = {cat: 0 for cat in categories}
    if "Subdomain" in df.columns:
        subdomain_counts = df["Subdomain"].value_counts()
        for subdomain, count in subdomain_counts.items():
            # Match subdomain to category
            matched = False
            for cat, aliases in categories.items():
                if subdomain in aliases:
                    counts[cat] += count
                    matched = True
                    break
            if not matched:
                logging.warning(f"  ⚠ Unrecognized subdomain for subtitle: '{subdomain}'")
    
    # Construct new subtitle string using dynamic category order
    subtitle_parts = []
    for cat in category_order:
        count = counts.get(cat, 0)
        # Use consistent separator for all categories
        subtitle_parts.append(f"{cat} – {count}")
    
    subtitle_text = "; ".join(subtitle_parts) + ";"
    
    # Find the shape to update - look for any subtitle text with category patterns
    # Strategy: Find textbox that contains subtitle-like patterns
    target_shape = None
    
    # First, try to find shape with existing subtitle pattern (contains category names)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        # Check if this looks like a subtitle (contains multiple category keywords)
        category_matches = sum(1 for cat in categories.keys() if cat in text)
        if category_matches >= 2:  # If it contains at least 2 category names, it's likely the subtitle
            target_shape = shape
            break
    
    if target_shape:
        logging.info(f"  Updating subtitle: {subtitle_text}")
        text_frame = target_shape.text_frame
        
        if len(text_frame.paragraphs) > 0:
            # Keep first paragraph (Title)
            # Remove subsequent paragraphs using XML
            text_body = text_frame._txBody
            p_elements = list(text_body.p_lst)
            for i in range(len(p_elements) - 1, 0, -1):  # Iterate backwards, stop before 0
                text_body.remove(p_elements[i])
            
            # Add new paragraph for subtitle
            p = text_frame.add_paragraph()
            p.text = subtitle_text
            p.font.size = Pt(14)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(255, 255, 255)  # White color
            p.space_before = Pt(6)
        else:
            logging.warning("  ⚠ Target shape has no paragraphs")
    else:
        logging.warning(f"  ⚠ Could not find subtitle shape with category patterns")
