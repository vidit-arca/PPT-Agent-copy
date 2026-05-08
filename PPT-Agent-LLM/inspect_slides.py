from pptx import Presentation
import glob

def list_slides():
    ppt_files = glob.glob("*.pptx")
    # Filter out temporary files
    ppt_files = [f for f in ppt_files if not f.startswith("~$") and "filled" not in f]
    
    if not ppt_files:
        print("No template found")
        return

    path = ppt_files[0]
    print(f"Analyzing: {path}")
    
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        else:
            # Try to find first text box
            for shape in slide.shapes:
                if shape.has_text_frame:
                    title = shape.text_frame.text
                    break
        print(f"Slide {i} (Index {i}): {title.strip()}")

if __name__ == "__main__":
    list_slides()
