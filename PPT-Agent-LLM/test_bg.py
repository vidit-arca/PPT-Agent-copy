from pptx import Presentation

prs = Presentation("tejomaya_filled_20260822_122640.pptx")
print(f"Total slides: {len(prs.slides)}")

for i, slide in enumerate(prs.slides):
    # print all shapes and their positions/sizes for slide 4 (index 3) and slide 5 (index 4)
    # the SEBI slides might be index 4, 5, etc.
    print(f"\n--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        print(f"  Shape {j}: {shape.shape_type} ({shape.name}) - Top: {shape.top}, Left: {shape.left}, Width: {shape.width}, Height: {shape.height}")
