import unittest
import sys
import os
from unittest.mock import MagicMock, patch
import pandas as pd

# Add src to path
sys.path.append(os.path.abspath(os.path.join(os.path.dirname(__file__), '..')))

from src.core.agent import paginate_and_fill
from src.services.ppt.layout import LayoutEngine

class TestPagination(unittest.TestCase):
    def setUp(self):
        self.prs = MagicMock()
        self.prs.slide_height.inches = 7.5
        self.prs.slide_width.inches = 13.33
        
        self.slide = MagicMock()
        self.prs.slides = MagicMock()
        self.prs.slides.__getitem__.return_value = self.slide
        self.prs.slides.__len__.return_value = 1
        self.prs.slides.index.return_value = 1
        
        self.shape = MagicMock()
        self.shape.has_table = True
        self.shape.height.inches = 2.0 # Header height
        self.shape.top = 1000 # This is raw units, let's use Inches for consistency if needed, but the code compares with Inches.
        # Actually, python-pptx uses Emu (int). Inches returns Emu.
        # But in my test I set self.shape.top = 1000.
        # And Inches(0.1) is 91440.
        # So 1000 - 1000 is 0 < 91440.
        # The problem is s.top - shape.top returns MagicMock if they are MagicMocks.
        # I need to ensure they are integers.
        from pptx.util import Inches
        self.shape.top = Inches(1)
        self.shape.left = Inches(1)
        
        self.table = MagicMock()
        self.shape.table = self.table
        self.table.rows = [MagicMock()] # Header row
        self.table.columns = [MagicMock(), MagicMock()]
        self.table.columns[0].width.inches = 2.0
        self.table.columns[1].width.inches = 8.0
        
        self.slide.shapes = [self.shape]
        
        # Mock LayoutEngine to control overflow
        self.layout_engine = MagicMock()
        self.layout_engine.check_horizontal_overflow.return_value = False
        self.layout_engine.estimate_row_height.return_value = 0.5
        self.layout_engine.analyze_split_point.return_value = 3
        
        # Mock duplicate_slide
        self.new_slide = MagicMock()
        self.new_slide.shapes = [self.shape] # Reuse shape for simplicity
        
    @patch('src.core.agent.duplicate_slide')
    @patch('src.core.agent.add_table_rows')
    @patch('src.core.agent.set_cell_text')
    @patch('src.core.agent.get_shapes_below')
    def test_pagination_overflow(self, mock_get_shapes, mock_set_cell, mock_add_rows, mock_dup_slide):
        # Setup
        mock_dup_slide.return_value = self.new_slide
        mock_get_shapes.return_value = []
        
        # Data: 10 rows
        df = pd.DataFrame({'Col1': [f'Val{i}' for i in range(10)]})
        
        # Layout: Max 3 rows per slide (Header=2.0, Row=0.5, Max=7.5, Margin=1.0 -> Available=6.5. 2.0+3*0.5=3.5 < 6.5. Wait.)
        # Let's force overflow after 3 rows
        # check_overflow(current_height, row_height, ...)
        
        def side_effect_overflow(curr, row, max_h, bottom_margin):
            # Simulate overflow if we have added 3 rows (current height starts at 2.0, +3*0.5 = 3.5)
            # Let's say max capacity is 3.5
            return curr + row > 3.6 
            
        self.layout_engine.check_overflow.side_effect = side_effect_overflow
        
        headers_map = {'Header1': ['Col1']}
        header_texts = ['Header1']
        header_to_colspec = {'Header1': ['Col1']}
        
        # Run
        paginate_and_fill(
            self.prs, 0, 0, df, 
            headers_map, header_texts, header_to_colspec, {}, 
            layout_engine=self.layout_engine
        )
        
        # Assert
        # Should have filled 3 rows on first slide
        # Should have duplicated slide
        mock_dup_slide.assert_called()
        
        print("Test finished successfully")

if __name__ == '__main__':
    unittest.main()
