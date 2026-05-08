from pptx import Presentation

# Analyze what's REALLY in the template
prs = Presentation('/Users/apple/Desktop/Akshayam/PPT-Agent/excel-to-ppt-service/templates/Akshayam Tejomaya edition.pptx')

print("=" * 70)
print("COMPLETE TEMPLATE ANALYSIS")
print("=" * 70)

for slide_idx, slide in enumerate(prs.slides):
    print(f"\n{'='*70}")
    print(f"SLIDE {slide_idx + 1}")
    print(f"{'='*70}")
    
    text_count = 0
    
    for shape_idx, shape in enumerate(slide.shapes):
        
        # Check if it has text
        has_text_frame = hasattr(shape, "text_frame")
        
        if has_text_frame and shape.text_frame:
            text = shape.text_frame.text
            
            if text.strip():  # Only show non-empty
                text_count += 1
                print(f"\n  [{text_count}] Shape: {shape.name}")
                print(f"      Type: {type(shape).__name__}")
                print(f"      Text: '{text}'")
                print(f"      Length: {len(text)} chars")
    
    if text_count == 0:
        print("  (No text content in this slide)")

print(f"\n{'='*70}")
print("Analysis complete!")
print("Look for text containing 'GistThereof', 'Summary', 'Date', etc.")
print(f"{'='*70}")
