import argparse
import pandas as pd
from pptx import Presentation


# ------------------------------
# SAFE TABLE FILLER (Option B)
# ------------------------------
def fill_table(table, df):
    """
    Fills a table WITHOUT adding rows.
    If data > table rows → only fills available rows safely.
    """

    max_rows = len(table.rows) - 1  # row 0 = header
    max_fill = min(max_rows, len(df))

    print(f"  Table has {max_rows} data rows, filling {max_fill}")

    for i in range(max_fill):
        row = table.rows[i + 1]  # skip header
        data = df.iloc[i]

        # Column 0 → S. No
        row.cells[0].text = str(i + 1)

        # Column 1 → Date (Month + Year)
        month = str(data.get("Month", "") or "")
        year = str(data.get("Year", "") or "")
        row.cells[1].text = f"{month} {year}".strip()

        # Column 2 → SubCategory
        row.cells[2].text = str(data.get("SubCategory", "") or "")

        # Column 3 → File Name
        row.cells[3].text = str(data.get("File Name", "") or "")

        # Column 4 → Summary
        row.cells[4].text = str(data.get("Summary", "") or "")

        print(f"    ✓ Filled row {i + 1}")

    if len(df) > max_rows:
        print(f"  ⚠ {len(df) - max_rows} rows could not fit (table too small)")


# ------------------------------
# MAIN AGENT
# ------------------------------
def fill_tejomaya(excel_path, template_path, output_path, sheet, section_map):
    print("\n===========================================================")
    print("TEJOMAYA TABLE FILLER — OPTION B (SAFE)")
    print("===========================================================\n")

    # Load Excel
    df = pd.read_excel(excel_path, sheet_name=sheet)
    df.columns = df.columns.str.strip()
    print(f"Loaded {len(df)} rows.")

    prs = Presentation(template_path)

    # Process each mapping
    for vertical, slide_idx in section_map.items():
        print(f"\nProcessing vertical '{vertical}' → Slide {slide_idx + 1}")

        df_v = df[df["Verticals"] == vertical]
        if df_v.empty:
            print(f"  ⚠ No rows for vertical {vertical}")
            continue

        slide = prs.slides[slide_idx]

        # FIND FIRST TABLE IN SLIDE
        table = None
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                break

        if table is None:
            print("  ❌ No table found on slide")
            continue

        # Fill table
        fill_table(table, df_v)

    # Save
    prs.save(output_path)
    print(f"\nSaved: {output_path}")
    print("===========================================================\n")


# ------------------------------
# CLI
# ------------------------------
if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--excel", required=True)
    parser.add_argument("--template", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument("--sheet", required=True)
    parser.add_argument("--map", nargs="+", required=True)
    args = parser.parse_args()

    # Parse mapping: "Companies Act:3" → {"Companies Act": 3}
    section_map = {}
    for m in args.map:
        key, val = m.split(":")
        section_map[key] = int(val)

    fill_tejomaya(
        args.excel,
        args.template,
        args.out,
        args.sheet,
        section_map
    )
